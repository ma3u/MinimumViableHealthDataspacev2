#!/usr/bin/env python3
"""
Generate a 10-slide PowerPoint for the Spanish Ministry of Health.

Audience: non-technical decision makers. No acronyms, no jargon, no
implementation details — just plain words, simple shapes, and people.

Output: docs/demos/spain-ehds-ministry-deck.pptx

Run:
    /tmp/pptx-venv/bin/python scripts/generate-spain-ehds-deck.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# ── Brand palette ─────────────────────────────────────────────────────────
NAVY   = RGBColor(0x14, 0x3D, 0x59)   # body headings
RED    = RGBColor(0xC3, 0x1E, 0x2E)   # Spain accent (warm red, not too bright)
GOLD   = RGBColor(0xE6, 0xA8, 0x17)   # Spain accent
TEAL   = RGBColor(0x2A, 0x9D, 0x8F)   # health / care
GREEN  = RGBColor(0x52, 0x8B, 0x55)   # research
GREY   = RGBColor(0x4F, 0x55, 0x60)   # body text
LIGHT  = RGBColor(0xF4, 0xEF, 0xE6)   # warm card background
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x1A, 0x1A, 0x1A)

# 16:9 = 13.333 × 7.5 in
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=GREY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, x, y, w, h, items, *, size=18, color=GREY,
                bold_first=False, line_spacing=1.15):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = "•  " + item
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        if bold_first and i == 0:
            run.font.bold = True
    return box


def add_box(slide, x, y, w, h, *, fill=LIGHT, line=None, line_w=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    s.text = ""
    return s


def add_circle(slide, cx, cy, r, *, fill=NAVY, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    return s


def add_arrow(slide, x1, y1, x2, y2, *, color=NAVY, weight=2.0):
    line = slide.shapes.add_connector(2, x1, y1, x2, y2)  # straight
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    # End arrow
    from pptx.oxml.ns import qn
    from lxml import etree
    ln = line.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")
    return line


def slide_header(slide, title, subtitle=None, accent=NAVY):
    # Top accent strip
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   0, 0, SLIDE_W, Inches(0.18))
    strip.fill.solid(); strip.fill.fore_color.rgb = accent
    strip.line.fill.background()
    add_text(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.7),
             title, size=30, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.05), Inches(12), Inches(0.45),
                 subtitle, size=16, color=GREY)
    # Bottom footer
    add_text(slide, Inches(0.6), Inches(7.05), Inches(12), Inches(0.3),
             "European Health Data Space  ·  A federated approach for Spain",
             size=10, color=GREY, align=PP_ALIGN.LEFT)


def add_person(slide, cx, cy, *, color=NAVY, label=None, label_size=14):
    """Simple stick-figure-ish person made from a circle + rounded rectangle."""
    head_r = Inches(0.22)
    body_w, body_h = Inches(0.55), Inches(0.55)
    add_circle(slide, cx, cy - Inches(0.35), head_r, fill=color)
    # Body / shoulders
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        cx - body_w / 2, cy - Inches(0.10), body_w, body_h,
    )
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.adjustments[0] = 0.3
    if label:
        add_text(slide, cx - Inches(1.2), cy + Inches(0.6),
                 Inches(2.4), Inches(0.4),
                 label, size=label_size, bold=True, color=color,
                 align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
# Solid navy banner top half
banner = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                            0, 0, SLIDE_W, Inches(4.2))
banner.fill.solid(); banner.fill.fore_color.rgb = NAVY
banner.line.fill.background()
# Spain accent stripe
stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                            0, Inches(4.2), SLIDE_W, Inches(0.25))
stripe.fill.solid(); stripe.fill.fore_color.rgb = GOLD
stripe.line.fill.background()
stripe2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             0, Inches(4.45), SLIDE_W, Inches(0.08))
stripe2.fill.solid(); stripe2.fill.fore_color.rgb = RED
stripe2.line.fill.background()

add_text(s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(1.3),
         "A Connected Health System for Spain",
         size=46, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(s, Inches(0.8), Inches(2.5), Inches(11.7), Inches(0.9),
         "Building a federated European Health Data Space",
         size=26, color=WHITE)
add_text(s, Inches(0.8), Inches(3.2), Inches(11.7), Inches(0.6),
         "across the 17 regions of Spain",
         size=22, color=GOLD, bold=True)

add_text(s, Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.5),
         "Better care for every citizen.  Better research for every region.",
         size=20, color=NAVY)
add_text(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.45),
         "Briefing for the Spanish Ministry of Health",
         size=16, color=GREY)
add_text(s, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.4),
         "May 2026",
         size=14, color=GREY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Challenge
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "The challenge today",
             "Spain's strength is regional autonomy. The cost is fragmentation.")

# Three problem cards
card_w = Inches(4.0)
card_h = Inches(4.6)
card_y = Inches(1.8)
card_x_offsets = [Inches(0.6), Inches(4.7), Inches(8.8)]

problems = [
    ("Patients", "A citizen who lives in Madrid and falls ill in Seville cannot share their medical history easily with the local doctor.", RED),
    ("Doctors", "Hospitals in different regions cannot see what care a patient already received elsewhere — tests get repeated, time is lost.", GOLD),
    ("Researchers", "Studying a rare disease across Spain means asking 17 separate regions for permission, data, and approvals.", TEAL),
]
for off, (h_text, body, accent) in zip(card_x_offsets, problems):
    add_box(s, off, card_y, card_w, card_h, fill=LIGHT)
    # accent bar on top
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             off, card_y, card_w, Inches(0.18))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_text(s, off + Inches(0.3), card_y + Inches(0.4),
             card_w - Inches(0.6), Inches(0.6),
             h_text, size=22, bold=True, color=NAVY)
    add_text(s, off + Inches(0.3), card_y + Inches(1.2),
             card_w - Inches(0.6), card_h - Inches(1.5),
             body, size=16, color=GREY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Vision
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "The vision",
             "One European Health Data Space.  Each region keeps its data and its rules.")

# Centre: large title circle
add_circle(s, Inches(6.67), Inches(4.0), Inches(1.0), fill=NAVY)
add_text(s, Inches(5.27), Inches(3.7), Inches(2.8), Inches(0.6),
         "Spain", size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Six regions around the centre — cardinal positions
regions = [
    ("Galicia",     Inches(2.5),  Inches(2.7)),
    ("Catalonia",   Inches(10.8), Inches(2.7)),
    ("Madrid",      Inches(6.67), Inches(2.0)),
    ("Andalusia",   Inches(6.67), Inches(6.1)),
    ("Basque Country", Inches(2.5),  Inches(5.4)),
    ("Valencia",    Inches(10.8), Inches(5.4)),
]
for name, cx, cy in regions:
    add_circle(s, cx, cy, Inches(0.55), fill=TEAL, line=NAVY)
    add_text(s, cx - Inches(1.2), cy + Inches(0.65),
             Inches(2.4), Inches(0.4),
             name, size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    # connector line back to Spain
    add_arrow(s, cx, cy, Inches(6.67), Inches(4.0), color=GREY, weight=1.0)

# Right side principles
add_text(s, Inches(8.0), Inches(6.4), Inches(5.0), Inches(0.4),
         "  +  16 more regions",
         size=12, color=GREY, align=PP_ALIGN.LEFT)

# Bottom: three short principles
y = Inches(6.7)
add_text(s, Inches(0.6), y, Inches(12), Inches(0.4),
         "Each region keeps its data.   Citizens carry one identity.   "
         "Trusted rules let information flow when it should.",
         size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Two ways data helps people
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "Two ways data helps people",
             "Both are equally important.  Both follow strict rules.")

# Two big cards
left_x  = Inches(0.6)
right_x = Inches(7.0)
card_y  = Inches(1.7)
card_w  = Inches(5.7)
card_h  = Inches(5.0)

# LEFT: Direct care
add_box(s, left_x, card_y, card_w, card_h, fill=LIGHT)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, card_y, card_w, Inches(0.4))
bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()
add_text(s, left_x + Inches(0.3), card_y + Inches(0.55),
         card_w - Inches(0.6), Inches(0.6),
         "1.  For your care",
         size=26, bold=True, color=TEAL)
add_text(s, left_x + Inches(0.3), card_y + Inches(1.2),
         card_w - Inches(0.6), Inches(0.5),
         "Doctors see what they need to treat you.",
         size=16, bold=True, color=NAVY)
add_bullets(s, left_x + Inches(0.3), card_y + Inches(1.9),
            card_w - Inches(0.6), Inches(2.6), [
    "Your medical record follows you across hospitals and regions.",
    "Emergency doctors can see allergies and prescriptions in seconds.",
    "Less repeated tests, less lost paperwork, faster care.",
    "You decide what is shared and you can withdraw consent.",
], size=15, color=GREY)

# RIGHT: Research and policy
add_box(s, right_x, card_y, card_w, card_h, fill=LIGHT)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, card_y, card_w, Inches(0.4))
bar.fill.solid(); bar.fill.fore_color.rgb = GREEN; bar.line.fill.background()
add_text(s, right_x + Inches(0.3), card_y + Inches(0.55),
         card_w - Inches(0.6), Inches(0.6),
         "2.  For better health policy",
         size=26, bold=True, color=GREEN)
add_text(s, right_x + Inches(0.3), card_y + Inches(1.2),
         card_w - Inches(0.6), Inches(0.5),
         "Researchers learn from millions of records.",
         size=16, bold=True, color=NAVY)
add_bullets(s, right_x + Inches(0.3), card_y + Inches(1.9),
            card_w - Inches(0.6), Inches(2.6), [
    "Hospitals and regions share patterns, never names.",
    "Researchers spot diseases earlier and find better treatments.",
    "The Ministry sees what is working in different regions.",
    "Every research request is approved and supervised.",
], size=15, color=GREY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — The five people in the system
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "The five people in the system",
             "Real roles you will meet in the live demonstration.")

people = [
    ("Citizen",     "Their health,\ntheir record.",                   TEAL),
    ("Hospital",    "Provides care,\nholds the records.",             NAVY),
    ("Researcher",  "Studies patterns\nto improve health.",            GREEN),
    ("Health Data Authority", "Approves and\nsupervises every request.", RED),
    ("Trust Centre", "Protects identity,\nresolves only when legal.",   GOLD),
]
n = len(people)
slot_w = Inches(2.4)
gap    = Inches(0.18)
total  = slot_w * n + gap * (n - 1)
start_x = (SLIDE_W - total) / 2
y = Inches(2.0)

for i, (name, sub, color) in enumerate(people):
    cx = start_x + slot_w * i + slot_w / 2 + (gap * i)
    add_circle(s, cx, y + Inches(0.9), Inches(0.85), fill=color)
    # body
    body = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        cx - Inches(0.85), y + Inches(1.6), Inches(1.7), Inches(1.0),
    )
    body.fill.solid(); body.fill.fore_color.rgb = color
    body.line.fill.background()
    body.adjustments[0] = 0.4
    # name
    add_text(s, cx - Inches(1.3), y + Inches(2.85),
             Inches(2.6), Inches(0.5),
             name, size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    # role
    add_text(s, cx - Inches(1.3), y + Inches(3.4),
             Inches(2.6), Inches(1.0),
             sub, size=13, color=GREY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — How it works for patients (primary use)
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "How it works for a citizen",
             "A patient from Madrid is on holiday in Valencia and feels unwell.",
             accent=TEAL)

# Three-step horizontal flow
step_y = Inches(2.4)
step_h = Inches(2.2)
step_w = Inches(3.6)
gap_x  = Inches(0.45)
total_w = step_w * 3 + gap_x * 2
start_x = (SLIDE_W - total_w) / 2

steps = [
    ("1.  Show consent",
     "The patient signs in with their national identity at the Valencian hospital.",
     TEAL),
    ("2.  Records arrive",
     "The Madrid hospital is asked, the rules check out, the record is delivered.",
     NAVY),
    ("3.  Care happens",
     "Doctor sees allergies, current medications, recent tests.  Treats safely.",
     GREEN),
]
for i, (title, body, color) in enumerate(steps):
    x = start_x + (step_w + gap_x) * i
    add_box(s, x, step_y, step_w, step_h, fill=LIGHT)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, step_y, step_w, Inches(0.18))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    add_text(s, x + Inches(0.3), step_y + Inches(0.32),
             step_w - Inches(0.6), Inches(0.5),
             title, size=20, bold=True, color=color)
    add_text(s, x + Inches(0.3), step_y + Inches(1.0),
             step_w - Inches(0.6), step_h - Inches(1.2),
             body, size=14, color=GREY)
    if i < 2:
        # arrow between cards
        ax1 = x + step_w + Inches(0.05)
        ax2 = x + step_w + gap_x - Inches(0.05)
        ay  = step_y + step_h / 2
        add_arrow(s, ax1, ay, ax2, ay, color=NAVY, weight=2.2)

# Bottom callout
add_text(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(0.45),
         "Three rules always apply",
         size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(5.9), Inches(12.1), Inches(0.4),
         "Citizen consent.    Local rules in both regions.    "
         "Every access logged for the citizen to see.",
         size=15, color=GREY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — How it works for research (secondary use)
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "How it works for research",
             "A researcher studies how diabetes treatment differs across regions.",
             accent=GREEN)

# Five-step flow
step_y = Inches(2.0)
step_w = Inches(2.4)
step_h = Inches(2.6)
gap_x  = Inches(0.18)
total_w = step_w * 5 + gap_x * 4
start_x = (SLIDE_W - total_w) / 2

steps = [
    ("Apply",  "Researcher submits a clear question and a study plan.",  GREEN),
    ("Approve","Health Data Authority checks purpose, ethics, scope.",   RED),
    ("Protect","Trust Centre replaces names with safe pseudonyms.",      GOLD),
    ("Analyse","Researcher works inside a secure room.  Data never leaves.", NAVY),
    ("Publish","Only summary results come out.  No individual is identifiable.", TEAL),
]
for i, (title, body, color) in enumerate(steps):
    x = start_x + (step_w + gap_x) * i
    add_box(s, x, step_y, step_w, step_h, fill=LIGHT)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, step_y, step_w, Inches(0.18))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    # number badge
    add_circle(s, x + Inches(0.4), step_y + Inches(0.5), Inches(0.22), fill=color)
    add_text(s, x + Inches(0.18), step_y + Inches(0.32),
             Inches(0.45), Inches(0.4),
             str(i + 1), size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.78), step_y + Inches(0.35),
             step_w - Inches(0.9), Inches(0.5),
             title, size=18, bold=True, color=color)
    add_text(s, x + Inches(0.25), step_y + Inches(1.0),
             step_w - Inches(0.5), step_h - Inches(1.2),
             body, size=12, color=GREY)

add_text(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(0.45),
         "What never happens",
         size=18, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.0),
         "Identifiable patient information never leaves the secure room.   "
         "No researcher sees a name, a national identity number, or a home address.   "
         "If a safety concern is found, only the Trust Centre — under legal authority — can re-link the data to a person.",
         size=14, color=GREY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — The Authority's role
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "The role of the Health Data Authority",
             "The Ministry's data office is the trusted gatekeeper for every secondary use.",
             accent=RED)

duties = [
    ("Approve research requests",
     "Reviews each application: purpose, lawful basis, proportionality, ethics."),
    ("Run the secure room",
     "Operates the protected environment where researchers analyse the data."),
    ("Publish a transparent register",
     "Every approval, every rejection, every research outcome is public."),
    ("Cooperate across borders",
     "Works with peer authorities in Germany, France, and the wider Union."),
    ("Enforce the rules",
     "Suspends or revokes permits and applies penalties for misuse."),
    ("Maintain the catalogue",
     "Publishes a public list of available datasets with quality labels."),
]
# 2 columns × 3 rows
col_w = Inches(6.0)
row_h = Inches(1.5)
gap = Inches(0.25)
start_x = Inches(0.6)
start_y = Inches(1.85)

for i, (h, body) in enumerate(duties):
    col = i % 2
    row = i // 2
    x = start_x + (col_w + gap) * col
    y = start_y + (row_h + gap) * row
    add_box(s, x, y, col_w, row_h, fill=LIGHT)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.18), row_h)
    bar.fill.solid(); bar.fill.fore_color.rgb = RED; bar.line.fill.background()
    add_text(s, x + Inches(0.4), y + Inches(0.18),
             col_w - Inches(0.6), Inches(0.45),
             h, size=17, bold=True, color=NAVY)
    add_text(s, x + Inches(0.4), y + Inches(0.7),
             col_w - Inches(0.6), Inches(0.7),
             body, size=13, color=GREY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Trust by design
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "Trust by design",
             "Citizens stay in control.  Every step is visible.",
             accent=GOLD)

trusts = [
    ("Citizen consent",
     "The citizen always says what may be shared, with whom, and for how long.  Consent can be withdrawn at any time."),
    ("Pseudonyms by default",
     "Identifiers are replaced with safe codes before any researcher sees the data.  The link can be re-made only when the law requires it."),
    ("A complete diary",
     "Every access — who, when, what, why — is recorded.  The citizen can read their own access diary."),
    ("Right to ask, change, erase",
     "Citizens can request their record, ask for corrections, and ask for erasure under European data protection law."),
]
y_top = Inches(1.95)
card_w = Inches(6.0)
card_h = Inches(2.4)
positions = [
    (Inches(0.6),  y_top),
    (Inches(6.85), y_top),
    (Inches(0.6),  y_top + card_h + Inches(0.3)),
    (Inches(6.85), y_top + card_h + Inches(0.3)),
]
for (x, y), (h, body) in zip(positions, trusts):
    add_box(s, x, y, card_w, card_h, fill=LIGHT)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.18), card_h)
    bar.fill.solid(); bar.fill.fore_color.rgb = GOLD; bar.line.fill.background()
    add_text(s, x + Inches(0.4), y + Inches(0.3),
             card_w - Inches(0.6), Inches(0.5),
             h, size=20, bold=True, color=NAVY)
    add_text(s, x + Inches(0.4), y + Inches(0.95),
             card_w - Inches(0.6), card_h - Inches(1.2),
             body, size=14, color=GREY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — A roadmap for Spain
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "A roadmap for Spain",
             "Three phases.  Each one delivers value on its own.",
             accent=NAVY)

phases = [
    ("Phase 1",  "Foundation",
     "Establish the Health Data Authority.  Agree common rules with the regions.  Set up the secure room.",
     TEAL),
    ("Phase 2",  "Connect three pioneer regions",
     "Madrid, Catalonia, Andalusia.  Citizens experience care that follows them.  First research approvals issued.",
     GREEN),
    ("Phase 3",  "Federate all 17 regions and join Europe",
     "Every region connected.  Cross-border with Germany, France, the wider Union.  Spain becomes a leader in trusted health data.",
     NAVY),
]
y = Inches(2.0)
phase_h = Inches(1.55)
gap_y = Inches(0.25)
for i, (label, title, body, color) in enumerate(phases):
    yy = y + (phase_h + gap_y) * i
    add_box(s, Inches(0.6), yy, Inches(12.1), phase_h, fill=LIGHT)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.6), yy, Inches(0.18), phase_h)
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    # phase label badge
    add_box(s, Inches(0.95), yy + Inches(0.32),
            Inches(1.4), Inches(0.55), fill=color)
    add_text(s, Inches(0.95), yy + Inches(0.32),
             Inches(1.4), Inches(0.55),
             label, size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(2.6), yy + Inches(0.18),
             Inches(10), Inches(0.55),
             title, size=22, bold=True, color=NAVY)
    add_text(s, Inches(2.6), yy + Inches(0.78),
             Inches(10), Inches(0.7),
             body, size=14, color=GREY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — The conversation we want to start
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "The conversation we want to start",
             "Working together with the Ministry and the regions.",
             accent=NAVY)

# Two columns: questions / what we offer
left_x = Inches(0.6); left_w = Inches(6.0)
right_x = Inches(6.85); right_w = Inches(6.0)

add_box(s, left_x, Inches(1.85), left_w, Inches(5.0), fill=LIGHT)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                         left_x, Inches(1.85), left_w, Inches(0.4))
bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
add_text(s, left_x + Inches(0.3), Inches(1.95),
         left_w - Inches(0.6), Inches(0.4),
         "Questions for the Ministry",
         size=20, bold=True, color=WHITE)
add_bullets(s, left_x + Inches(0.3), Inches(2.55),
            left_w - Inches(0.6), Inches(4.2), [
    "Which regions would be the first three to connect?",
    "How will the Health Data Authority be staffed and funded?",
    "Which existing public health programmes benefit most quickly?",
    "How should citizens be informed and consulted?",
    "Which European partners should Spain federate with first?",
], size=15, color=GREY)

add_box(s, right_x, Inches(1.85), right_w, Inches(5.0), fill=LIGHT)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                         right_x, Inches(1.85), right_w, Inches(0.4))
bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()
add_text(s, right_x + Inches(0.3), Inches(1.95),
         right_w - Inches(0.6), Inches(0.4),
         "What we bring to the table",
         size=20, bold=True, color=WHITE)
add_bullets(s, right_x + Inches(0.3), Inches(2.55),
            right_w - Inches(0.6), Inches(4.2), [
    "A working live demonstration with realistic data.",
    "Experience federating data across European regions.",
    "Templates for the legal, organisational, and operational setup.",
    "A roadmap aligned with the European timetable.",
    "A team that has done this before, ready to support each region.",
], size=15, color=GREY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Thank you / contact
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
banner = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(4.0))
banner.fill.solid(); banner.fill.fore_color.rgb = NAVY; banner.line.fill.background()
stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                            0, Inches(4.0), SLIDE_W, Inches(0.25))
stripe.fill.solid(); stripe.fill.fore_color.rgb = GOLD; stripe.line.fill.background()
stripe2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             0, Inches(4.25), SLIDE_W, Inches(0.08))
stripe2.fill.solid(); stripe2.fill.fore_color.rgb = RED; stripe2.line.fill.background()

add_text(s, Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.2),
         "Thank you",
         size=54, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(2.55), Inches(11.7), Inches(0.7),
         "Better care for every citizen.  Better research for every region.",
         size=22, color=WHITE)

add_text(s, Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.5),
         "Live demonstration available on request",
         size=18, bold=True, color=NAVY)
add_text(s, Inches(0.8), Inches(5.55), Inches(11.7), Inches(0.4),
         "Walk through the five roles, the patient journey, and the research approval process — end to end.",
         size=14, color=GREY)
add_text(s, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.4),
         "Sopra Steria  ·  European Health Data Space practice",
         size=13, color=GREY)


# ── Save ──────────────────────────────────────────────────────────────────
out = Path(__file__).resolve().parents[1] / "docs" / "demos" / "spain-ehds-ministry-deck.pptx"
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(out)
print(f"wrote {out}  ({len(prs.slides)} slides)")
